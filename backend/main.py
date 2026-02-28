from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os
import tempfile
import base64
from typing import List, Optional
import requests
import io

app = FastAPI()

# Enable CORS for Next.js frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:3001"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class Slide(BaseModel):
    num: int
    type: str
    title: str
    content: str
    speakerNote: Optional[str] = ""
    mermaid: Optional[str] = None

class PresentationRequest(BaseModel):
    title: str
    slides: List[Slide]
    mode: str

def render_mermaid_to_image(mermaid_code: str, output_path: str):
    """Render Mermaid diagram to PNG using Mermaid.ink API"""
    try:
        # Encode mermaid code
        graphbiz = base64.b64encode(mermaid_code.encode('utf-8')).decode('ascii')
        
        # Use mermaid.ink free API service
        url = f"https://mermaid.ink/img/{graphbiz}?theme=dark&bgColor=!0a0a0f"
        
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        
        # Save image
        img = Image.open(io.BytesIO(response.content))
        
        # Ensure dark background
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        img.save(output_path, 'PNG')
        return True
        
    except Exception as e:
        print(f"Error rendering mermaid: {e}")
        return False

@app.post("/generate-pptx")
async def generate_pptx(data: PresentationRequest):
    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Dark theme colors
        bg_color = RGBColor(10, 10, 15)  # #0a0a0f
        text_color = RGBColor(232, 232, 240)  # #e8e8f0
        accent_color = RGBColor(124, 58, 237)  # #7c3aed
        
        for slide_data in data.slides:
            # Choose layout based on slide type
            if slide_data.type == "title":
                layout = prs.slide_layouts[0]  # Title slide
            else:
                layout = prs.slide_layouts[1]  # Title and content
            
            slide = prs.slides.add_slide(layout)
            
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
            # Set title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_data.title
                title_frame = title_shape.text_frame
                for paragraph in title_frame.paragraphs:
                    paragraph.font.size = Pt(44) if slide_data.type == "title" else Pt(32)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = accent_color
            
            # Add content
            if slide_data.mermaid:
                # Render mermaid diagram as image
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as img_tmp:
                    img_path = img_tmp.name
                
                try:
                    render_mermaid_to_image(slide_data.mermaid, img_path)
                    
                    # Add image to slide
                    left = Inches(1)
                    top = Inches(2)
                    height = Inches(4.5)
                    slide.shapes.add_picture(img_path, left, top, height=height)
                    
                finally:
                    if os.path.exists(img_path):
                        os.unlink(img_path)
                        
            elif len(slide.shapes) > 1:
                content_shape = slide.shapes[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                # Parse content lines
                lines = slide_data.content.split('\n')
                for line in lines:
                    if line.strip():
                        p = text_frame.add_paragraph()
                        p.text = line.strip().replace('• ', '')
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.color.rgb = text_color
                        if line.strip().startswith('•'):
                            p.level = 1
            
            # Add speaker notes
            if slide_data.speakerNote:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data.speakerNote
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            prs.save(tmp.name)
            tmp_path = tmp.name
        
        # Return file
        return FileResponse(
            tmp_path,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=f'{data.title.replace(" ", "_")}.pptx',
            background=lambda: os.unlink(tmp_path)  # Clean up after sending
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/")
async def root():
    return {"message": "Research KISSer PowerPoint Generator API"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
